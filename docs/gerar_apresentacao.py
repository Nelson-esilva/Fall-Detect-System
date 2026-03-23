"""Gera apresentação PowerPoint do projeto Fall-Detect-System."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

OUTPUT = Path(__file__).parent / "Apresentacao_Fall_Detect_System.pptx"

DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x00, 0x96, 0xC7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
RED = RGBColor(0xE7, 0x4C, 0x3C)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
YELLOW = RGBColor(0xF1, 0xC4, 0x0F)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT_GRAY, bullet_color=ACCENT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(6)
        p.level = 0
    return tf


def add_colored_box(slide, left, top, width, height, fill_color, text="",
                    font_size=14, text_color=WHITE):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(font_size)
        tf.paragraphs[0].font.color.rgb = text_color
        tf.paragraphs[0].font.name = "Segoe UI"
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def slide_capa(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BG)
    add_colored_box(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT)
    add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "Sistema de Detecção de Quedas\ncom Inteligência Artificial",
             font_size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.6),
             "Deep Learning  •  ESP32  •  App Mobile  •  MQTT",
             font_size=20, color=ACCENT, align=PP_ALIGN.CENTER)
    add_colored_box(slide, Inches(3.5), Inches(4.2), Inches(6), Inches(0.02), ACCENT)
    add_text(slide, Inches(1), Inches(4.6), Inches(11), Inches(0.5),
             "Nelson Emeliano Silva", font_size=22, color=WHITE,
             bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.5),
             "Orientador: Prof. Angilberto Muniz Ferreira Sobrinho",
             font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
             "PAIC/FAPEAM — Universidade do Estado do Amazonas (UEA)",
             font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_colored_box(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), ACCENT)


def slide_padrao(prs, titulo, conteudo_fn):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_colored_box(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT)
    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             titulo, font_size=30, color=WHITE, bold=True)
    add_colored_box(slide, Inches(0.8), Inches(1.1), Inches(4), Inches(0.04), ACCENT)
    conteudo_fn(slide)
    add_colored_box(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), ACCENT)
    return slide


def slide_problema(prs):
    def content(slide):
        items = [
            "Quedas são a 2ª causa de morte acidental no mundo (OMS)",
            "Idosos e pessoas com mobilidade reduzida são os mais afetados",
            "Tempo de socorro é fator crítico para a sobrevivência",
            "Sistemas manuais (botões de pânico) falham quando a vítima está inconsciente",
            "Necessidade de detecção automática, em tempo real, sem intervenção do usuário",
        ]
        add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(7), Inches(4.5),
                        items, font_size=20)
        add_colored_box(slide, Inches(8.5), Inches(2), Inches(4), Inches(3),
                        RGBColor(0x22, 0x22, 0x3E))
        add_text(slide, Inches(8.7), Inches(2.3), Inches(3.6), Inches(0.5),
                 "O DESAFIO", font_size=16, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(8.7), Inches(2.9), Inches(3.6), Inches(2),
                 "Reduzir o tempo entre a queda e o socorro, "
                 "utilizando IA para detectar automaticamente e alertar cuidadores em segundos.",
                 font_size=15, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    slide_padrao(prs, "Problema e Motivação", content)


def slide_objetivo(prs):
    def content(slide):
        add_text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(1),
                 "Desenvolver um sistema completo de detecção de quedas que integra:",
                 font_size=20, color=LIGHT_GRAY)
        boxes = [
            (Inches(0.8), "Visão Computacional\n+ Deep Learning", "Câmera + modelo CNN-LSTM\nanalisa vídeo em tempo real"),
            (Inches(4.9), "Hardware IoT\n(ESP32)", "Alarme sonoro local\n(buzzer + LEDs) via MQTT/Serial"),
            (Inches(9.0), "App Mobile\n(React Native)", "Notificação remota ao cuidador\n(alarme + vibração + emergência)"),
        ]
        for left, title, desc in boxes:
            add_colored_box(slide, left, Inches(2.8), Inches(3.6), Inches(3.2),
                            RGBColor(0x22, 0x22, 0x3E))
            add_text(slide, left + Inches(0.2), Inches(3.0), Inches(3.2), Inches(1),
                     title, font_size=18, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
            add_text(slide, left + Inches(0.2), Inches(4.0), Inches(3.2), Inches(1.5),
                     desc, font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    slide_padrao(prs, "Objetivo do Projeto", content)


def slide_arquitetura(prs):
    def content(slide):
        components = [
            (Inches(0.5), Inches(2.8), "PC + Câmera\nMobileNetV2 + LSTM\nDetecção de Queda", RGBColor(0x1B, 0x4F, 0x72)),
            (Inches(4.6), Inches(2.8), "Broker MQTT\n(Mosquitto)\nporta 1883 / 9001 WS", RGBColor(0x4A, 0x23, 0x5A)),
            (Inches(8.7), Inches(1.8), "App Mobile\n(React Native)\nAlarme + Vibração", RGBColor(0x0E, 0x6B, 0x3A)),
            (Inches(8.7), Inches(4.2), "ESP32\nBuzzer + LEDs\nAlarme Local", RGBColor(0x7B, 0x24, 0x1C)),
        ]
        for left, top, text, color in components:
            add_colored_box(slide, left, top, Inches(3.5), Inches(1.8), color,
                            text, font_size=14, text_color=WHITE)
        # Setas (representadas com texto)
        add_text(slide, Inches(4.0), Inches(3.3), Inches(0.8), Inches(0.5),
                 ">>>", font_size=24, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(8.1), Inches(2.3), Inches(0.8), Inches(0.5),
                 ">>>", font_size=24, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(8.1), Inches(4.8), Inches(0.8), Inches(0.5),
                 ">>>", font_size=24, color=RED, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(3.6), Inches(3.9), Inches(1.5), Inches(0.4),
                 "publish", font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(7.8), Inches(1.8), Inches(1.5), Inches(0.4),
                 "subscribe", font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(7.8), Inches(5.4), Inches(1.5), Inches(0.4),
                 "subscribe", font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    slide_padrao(prs, "Arquitetura Geral do Sistema", content)


def slide_modelo_ia(prs):
    def content(slide):
        layers = [
            ("1. MobileNetV2 (CNN)", "Extrai características visuais de cada frame\nTransfer Learning — pesos ImageNet", RGBColor(0x1B, 0x4F, 0x72)),
            ("2. TimeDistributed", "Aplica a CNN nos 20 frames da janela\ntemporal individualmente", RGBColor(0x2C, 0x3E, 0x50)),
            ("3. LSTM (64 unidades)", "Analisa sequência temporal para\nidentificar padrão de queda", RGBColor(0x4A, 0x23, 0x5A)),
            ("4. Dense (sigmoid)", "Classificação binária:\nNormal (0) ou Queda (1)", RGBColor(0x7B, 0x24, 0x1C)),
        ]
        for i, (title, desc, color) in enumerate(layers):
            top = Inches(1.5) + Inches(i * 1.4)
            add_colored_box(slide, Inches(0.8), top, Inches(5.5), Inches(1.2), color)
            add_text(slide, Inches(1.0), top + Inches(0.1), Inches(2.5), Inches(0.5),
                     title, font_size=16, color=YELLOW, bold=True)
            add_text(slide, Inches(1.0), top + Inches(0.55), Inches(5), Inches(0.6),
                     desc, font_size=13, color=LIGHT_GRAY)
            if i < len(layers) - 1:
                add_text(slide, Inches(3.2), top + Inches(1.15), Inches(0.5), Inches(0.3),
                         "▼", font_size=16, color=ACCENT, align=PP_ALIGN.CENTER)
        add_colored_box(slide, Inches(7.2), Inches(1.5), Inches(5.3), Inches(5.5),
                        RGBColor(0x22, 0x22, 0x3E))
        add_text(slide, Inches(7.4), Inches(1.7), Inches(5), Inches(0.5),
                 "Entrada e Saída", font_size=18, color=ACCENT, bold=True)
        specs = [
            "Entrada: (batch, 20, 128, 128, 3)",
            "= 20 frames RGB de 128x128 pixels",
            "",
            "Saída: valor entre 0.0 e 1.0",
            "  > 0.75 por 3 frames → QUEDA",
            "  ≤ 0.50 → NORMAL",
            "",
            "Inferência: a cada 5 frames",
            "Confirmação temporal: 3 predições",
            "consecutivas acima do limiar",
        ]
        for i, line in enumerate(specs):
            add_text(slide, Inches(7.4), Inches(2.4) + Inches(i * 0.4), Inches(5), Inches(0.4),
                     line, font_size=14, color=LIGHT_GRAY if line else WHITE)
    slide_padrao(prs, "Modelo de IA — CNN + LSTM", content)


def slide_dataset(prs):
    def content(slide):
        add_text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
                 "UR Fall Detection Dataset — Universidade de Rzeszów (Polônia)",
                 font_size=20, color=ACCENT, bold=True)
        items = [
            "30 sequências de quedas simuladas",
            "40 sequências de atividades diárias normais (ADL)",
            "Câmeras RGB + sensores de profundidade + acelerômetro",
            "Resolução processada: 128×128 pixels, normalizada [0,1]",
        ]
        add_bullet_list(slide, Inches(0.8), Inches(2.2), Inches(6), Inches(2.5),
                        items, font_size=18)

        add_text(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(0.5),
                 "Pipeline ETL (Extração, Transformação, Carregamento)",
                 font_size=18, color=ACCENT, bold=True)
        steps = [
            ("Varredura\nRecursiva", RGBColor(0x1B, 0x4F, 0x72)),
            ("PNG → AVI\n(30 fps)", RGBColor(0x2C, 0x3E, 0x50)),
            ("Resize\n128×128", RGBColor(0x4A, 0x23, 0x5A)),
            ("Sliding Window\n(passo 10)", RGBColor(0x7B, 0x24, 0x1C)),
            ("Split 80/20\nEstratificado", RGBColor(0x0E, 0x6B, 0x3A)),
        ]
        for i, (text, color) in enumerate(steps):
            left = Inches(0.8) + Inches(i * 2.4)
            add_colored_box(slide, left, Inches(5.2), Inches(2.0), Inches(1.3),
                            color, text, font_size=13)
            if i < len(steps) - 1:
                add_text(slide, left + Inches(2.0), Inches(5.55), Inches(0.5), Inches(0.4),
                         "→", font_size=22, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    slide_padrao(prs, "Dataset e Pipeline de Dados", content)


def slide_treinamento(prs):
    def content(slide):
        col1 = [
            ("Transfer Learning", "MobileNetV2 pré-treinada (ImageNet)\nFine-tuning: 30 últimas camadas"),
            ("Regularização", "Dropout 0.5 na saída LSTM\nEarly Stopping (paciência = 5 épocas)"),
            ("Otimização", "Adam (lr = 1e-4)\nLoss: binary crossentropy"),
        ]
        for i, (title, desc) in enumerate(col1):
            top = Inches(1.5) + Inches(i * 1.8)
            add_colored_box(slide, Inches(0.8), top, Inches(5.5), Inches(1.5),
                            RGBColor(0x22, 0x22, 0x3E))
            add_text(slide, Inches(1.0), top + Inches(0.15), Inches(5), Inches(0.4),
                     title, font_size=17, color=ACCENT, bold=True)
            add_text(slide, Inches(1.0), top + Inches(0.6), Inches(5), Inches(0.8),
                     desc, font_size=14, color=LIGHT_GRAY)

        col2 = [
            ("Hiperparâmetros", "Batch size: 4\nÉpocas: 30\nSequência: 20 frames\nResolução: 128×128"),
            ("Conversão TFLite", "Quantização INT8\nReduz tamanho ~4x\nMais rápido em CPU"),
        ]
        for i, (title, desc) in enumerate(col2):
            top = Inches(1.5) + Inches(i * 2.5)
            add_colored_box(slide, Inches(7.2), top, Inches(5.3), Inches(2.2),
                            RGBColor(0x22, 0x22, 0x3E))
            add_text(slide, Inches(7.4), top + Inches(0.15), Inches(5), Inches(0.4),
                     title, font_size=17, color=ACCENT, bold=True)
            add_text(slide, Inches(7.4), top + Inches(0.6), Inches(5), Inches(1.5),
                     desc, font_size=14, color=LIGHT_GRAY)
    slide_padrao(prs, "Treinamento e Técnicas", content)


def slide_esp32(prs):
    def content(slide):
        add_text(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
                 "Comunicação PC ↔ ESP32", font_size=20, color=ACCENT, bold=True)
        modes = [
            ("Serial (USB)", "Simples, sem WiFi\nIdeal para testes", RGBColor(0x1B, 0x4F, 0x72)),
            ("MQTT (WiFi)", "Sem fio, escalável\nProdução", RGBColor(0x0E, 0x6B, 0x3A)),
        ]
        for i, (title, desc, color) in enumerate(modes):
            left = Inches(0.8) + Inches(i * 3.0)
            add_colored_box(slide, left, Inches(2.2), Inches(2.7), Inches(1.5), color)
            add_text(slide, left + Inches(0.15), Inches(2.3), Inches(2.4), Inches(0.4),
                     title, font_size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
            add_text(slide, left + Inches(0.15), Inches(2.8), Inches(2.4), Inches(0.8),
                     desc, font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

        add_text(slide, Inches(0.8), Inches(4.2), Inches(5.5), Inches(0.5),
                 "Componentes de Hardware", font_size=18, color=ACCENT, bold=True)
        hw = [
            "Buzzer passivo (GPIO 18) — alerta sonoro",
            "LED vermelho (GPIO 19) — alerta visual de queda",
            "LED verde (GPIO 21) — status do sistema",
            "Botão de teste (GPIO 0) — teste manual",
        ]
        add_bullet_list(slide, Inches(0.8), Inches(4.8), Inches(5.5), Inches(2.5),
                        hw, font_size=15)

        add_colored_box(slide, Inches(7.2), Inches(1.5), Inches(5.3), Inches(5.5),
                        RGBColor(0x22, 0x22, 0x3E))
        add_text(slide, Inches(7.4), Inches(1.7), Inches(5), Inches(0.5),
                 "Payload JSON", font_size=18, color=ACCENT, bold=True)
        json_text = (
            '{\n'
            '  "alert": "FALL_DETECTED",\n'
            '  "confidence": 0.95,\n'
            '  "timestamp": "2026-03-23...",\n'
            '  "metadata": {\n'
            '    "frame_id": 1234,\n'
            '    "model": "CNN-LSTM"\n'
            '  }\n'
            '}'
        )
        add_text(slide, Inches(7.6), Inches(2.4), Inches(4.8), Inches(4),
                 json_text, font_size=14, color=GREEN, font_name="Consolas")
    slide_padrao(prs, "Integração com Hardware — ESP32", content)


def slide_app_mobile(prs):
    def content(slide):
        add_text(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
                 "React Native + Expo SDK 54 (TypeScript)",
                 font_size=20, color=ACCENT, bold=True)
        features = [
            "Conexão MQTT via WebSocket (porta 9001)",
            "Alarme sonoro + vibração ao detectar queda",
            "Tela de alarme em tela cheia com botão de emergência",
            "Histórico de eventos com filtros por tipo",
            "Persistência local (AsyncStorage)",
            "Limiar de confiança configurável",
            "Discagem direta para emergência (SAMU 192)",
        ]
        add_bullet_list(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5),
                        features, font_size=16)

        screens = [
            ("Dashboard", "Status de conexão\nÚltimos eventos", RGBColor(0x1B, 0x4F, 0x72)),
            ("Alarme", "Tela cheia\nBotão emergência", RGBColor(0x7B, 0x24, 0x1C)),
            ("Histórico", "Lista filtrada\nDetalhes por evento", RGBColor(0x4A, 0x23, 0x5A)),
            ("Config", "IP do broker\nLimiar de confiança", RGBColor(0x0E, 0x6B, 0x3A)),
        ]
        for i, (title, desc, color) in enumerate(screens):
            left = Inches(7.2) + Inches(i % 2) * Inches(2.7)
            top = Inches(1.5) + Inches(i // 2) * Inches(2.8)
            add_colored_box(slide, left, top, Inches(2.4), Inches(2.5), color)
            add_text(slide, left + Inches(0.1), top + Inches(0.2), Inches(2.2), Inches(0.4),
                     title, font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
            add_text(slide, left + Inches(0.1), top + Inches(0.8), Inches(2.2), Inches(1.5),
                     desc, font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    slide_padrao(prs, "Aplicativo Mobile", content)


def slide_resultados(prs):
    def content(slide):
        results = [
            "Sistema funcional ponta a ponta (câmera → detecção → alerta)",
            "Modelo CNN-LSTM treinado e validado no UR Fall Detection Dataset",
            "Pipeline ETL automatizado (PNG → vídeo → amostras de treino)",
            "Detecção em tempo real via webcam ou arquivo de vídeo",
            "Integração simultânea: ESP32 (alarme local) + App (alarme remoto)",
            "Protocolo MQTT padronizado com payload JSON",
            "Conversão TFLite para inferência otimizada em CPU",
        ]
        add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(3),
                        results, font_size=18)

        add_text(slide, Inches(0.8), Inches(5.0), Inches(11), Inches(0.5),
                 "Stack Tecnológico", font_size=18, color=ACCENT, bold=True)
        stack = [
            ("IA/ML", "TensorFlow + Keras", RGBColor(0x1B, 0x4F, 0x72)),
            ("Visão", "OpenCV", RGBColor(0x2C, 0x3E, 0x50)),
            ("Backend", "Python 3.10", RGBColor(0x4A, 0x23, 0x5A)),
            ("Comunicação", "MQTT / Mosquitto", RGBColor(0x7B, 0x24, 0x1C)),
            ("Mobile", "React Native / Expo", RGBColor(0x0E, 0x6B, 0x3A)),
            ("Hardware", "ESP32 + Arduino", RGBColor(0x6C, 0x3D, 0x10)),
        ]
        for i, (label, tech, color) in enumerate(stack):
            left = Inches(0.8) + Inches(i * 2.05)
            add_colored_box(slide, left, Inches(5.5), Inches(1.85), Inches(1.4), color)
            add_text(slide, left + Inches(0.05), Inches(5.6), Inches(1.75), Inches(0.4),
                     label, font_size=12, color=YELLOW, bold=True, align=PP_ALIGN.CENTER)
            add_text(slide, left + Inches(0.05), Inches(6.05), Inches(1.75), Inches(0.6),
                     tech, font_size=13, color=WHITE, align=PP_ALIGN.CENTER)
    slide_padrao(prs, "Resultados Alcançados", content)


def slide_dificuldades(prs):
    def content(slide):
        challenges = [
            ("Compatibilidade TensorFlow / Windows",
             "TF 2.11+ não suporta GPU nativo no Windows.\nSolução: Python 3.10 + TF 2.10, inferência otimizada em CPU via TFLite."),
            ("Heterogeneidade do Dataset",
             "Estrutura de pastas inconsistente no UR Fall.\nSolução: busca recursiva (os.walk) agnóstica à estrutura."),
            ("Performance em Hardware Limitado",
             "Inferência pesada para CPU de notebook.\nSolução: skip frames, confirmação temporal, TFLite com quantização INT8."),
            ("Falsos Positivos Temporais",
             "Modelo alertava queda antes da pessoa cair de fato.\nSolução: limiar conservador (0.75) + 3 predições consecutivas."),
        ]
        for i, (title, desc) in enumerate(challenges):
            top = Inches(1.5) + Inches(i * 1.4)
            add_colored_box(slide, Inches(0.8), top, Inches(11.5), Inches(1.2),
                            RGBColor(0x22, 0x22, 0x3E))
            add_text(slide, Inches(1.0), top + Inches(0.1), Inches(11), Inches(0.4),
                     title, font_size=16, color=ORANGE, bold=True)
            add_text(slide, Inches(1.0), top + Inches(0.55), Inches(11), Inches(0.6),
                     desc, font_size=13, color=LIGHT_GRAY)
    slide_padrao(prs, "Dificuldades e Soluções", content)


def slide_conclusao(prs):
    def content(slide):
        add_text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
                 "Conclusão", font_size=22, color=ACCENT, bold=True)
        conclusao = [
            "Sistema completo e funcional: detecção → alerta local + remoto",
            "Arquitetura modular e extensível (novos sensores, mais dados)",
            "Viabilidade técnica e científica demonstrada para aplicação real",
            "Código aberto no GitHub para reprodutibilidade",
        ]
        add_bullet_list(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(2.5),
                        conclusao, font_size=18)

        add_text(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(0.5),
                 "Próximos Passos", font_size=22, color=ACCENT, bold=True)
        proximos = [
            "Testes em ambiente residencial real com idosos",
            "Ampliar dataset com mais vídeos e cenários diversos",
            "Otimização para deploy embarcado (Raspberry Pi / Jetson Nano)",
            "Publicação de artigo científico com resultados finais",
        ]
        add_bullet_list(slide, Inches(0.8), Inches(5.2), Inches(11), Inches(2),
                        proximos, font_size=18)
    slide_padrao(prs, "Conclusão e Próximos Passos", content)


def slide_referencias(prs):
    def content(slide):
        refs = [
            "KWOLEK, B.; KEPSKI, M. Human fall detection on embedded platform using depth maps "
            "and wireless accelerometer. Computer Methods and Programs in Biomedicine, v. 117, n. 3, "
            "p. 489–501, 2014.",
            "",
            "LU, N.; WU, Y.; FENG, L.; SONG, J. Deep Learning for Fall Detection: Three-Dimensional "
            "CNN Combined With LSTM on Video Kinematic Data. IEEE JBHI, v. 23, n. 1, p. 314–323, 2019.",
            "",
            "CHHETRI, S. et al. Deep learning for vision-based fall detection system: Enhanced optical "
            "dynamic flow. arXiv:2104.05744, 2021.",
            "",
            "MUBASHIR, M.; SHAO, L.; SEED, L. A survey on fall detection: Principles and approaches. "
            "Neurocomputing, v. 100, p. 144–152, 2013.",
            "",
            "QIAN, B.; LIU, L. DeepFall: Skeleton-based fall detection using recurrent neural networks. "
            "Master Thesis, KU Leuven, 2019.",
        ]
        y = Inches(1.5)
        for ref in refs:
            if ref:
                add_text(slide, Inches(0.8), y, Inches(11.5), Inches(0.6),
                         ref, font_size=12, color=LIGHT_GRAY)
                y += Inches(0.55)
            else:
                y += Inches(0.1)
    slide_padrao(prs, "Referências Bibliográficas", content)


def slide_obrigado(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_colored_box(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT)
    add_text(slide, Inches(1), Inches(2.0), Inches(11), Inches(1),
             "Obrigado!", font_size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_colored_box(slide, Inches(4.5), Inches(3.2), Inches(4), Inches(0.03), ACCENT)
    add_text(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.5),
             "Nelson Emeliano Silva", font_size=22, color=ACCENT,
             bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
             "github.com/Nelson-esilva/Fall-Detect-System",
             font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
             "PAIC/FAPEAM — UEA — 2025/2026",
             font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_colored_box(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), ACCENT)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_capa(prs)
    slide_problema(prs)
    slide_objetivo(prs)
    slide_arquitetura(prs)
    slide_modelo_ia(prs)
    slide_dataset(prs)
    slide_treinamento(prs)
    slide_esp32(prs)
    slide_app_mobile(prs)
    slide_resultados(prs)
    slide_dificuldades(prs)
    slide_conclusao(prs)
    slide_referencias(prs)
    slide_obrigado(prs)

    prs.save(str(OUTPUT))
    print(f"Apresentação salva em: {OUTPUT}")
    print(f"Total de slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
